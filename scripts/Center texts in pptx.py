
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

# Specify the folder path where the PowerPoint files are located
folder_path = r"C:\Projects\ConvertedPPTX White BG"

# Iterate through the PowerPoint files in the folder
for filename in os.listdir(folder_path):
    if filename.endswith('.pptx'):
        # Create a presentation object for the PowerPoint file
        presentation = Presentation(os.path.join(folder_path, filename))

        # Iterate through the slides in the presentation
        for slide in presentation.slides:
            # Iterate through the shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    # Center-align text in all text frames
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER  # Center-align text

        # Save the modified presentation
        presentation.save(os.path.join(folder_path, filename))

print("Text alignment updated to center in all PowerPoint files.")

