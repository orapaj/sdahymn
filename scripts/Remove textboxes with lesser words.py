from pptx import Presentation

# Open PowerPoint file
prs = Presentation(r"C:\Users\Admin\Desktop\Output pptx\001.pptx")  

# Get first slide 
slide = prs.slides[0]  

# Get list of textbox shapes
textboxes = [shp for shp in slide.shapes if shp.has_text_frame]

# Track max words  
max_words = 0
max_tb = None

# Loop through textboxes
for tb in textboxes:
  num_words = len(tb.text.split()) 
  if num_words > max_words:
    max_words = num_words
    max_tb = tb

for tb in textboxes:
  if tb != max_tb:  
    tb._element.getparent().remove(tb._element)

# Save modified file
prs.save('modified.pptx') 

print("Script complete!")
