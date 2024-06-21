import os
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Specify the folder path where the PowerPoint files are located
folder_path = r"C:\ConvertedPPTX copy"

# Iterate through the PowerPoint files in the folder
for filename in os.listdir(folder_path):
    if filename.endswith('.pptx'):
        # Create a presentation object for the PowerPoint file
        presentation = Presentation(os.path.join(folder_path, filename))

        # Iterate through the slides in the presentation
        for slide in presentation.slides:
            # Calculate the slide dimensions
            slide_width = presentation.slide_width
            slide_height = presentation.slide_height

            # Add a white rectangle (shape) to the slide
            left = top = Inches(0)  # Adjust the position as needed
            width = slide_width
            height = slide_height
            white_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )

            # Set the fill color of the white shape to white
            fill = white_shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)  # White color

            # Send the white shape to the back of the slide
            slide.shapes._spTree.remove(white_shape._element)
            slide.shapes._spTree.insert(2, white_shape._element)  # Place it at the back

        # Save the modified presentation
        presentation.save(os.path.join(folder_path, filename))

print("White shapes inserted and placed at the back in all PowerPoint files.")
