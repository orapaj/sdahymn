import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Specify the folder path where the PowerPoint files are located
folder_path = r"C:\Projects\ConvertedPPTX White BG"

# Iterate through the PowerPoint files in the folder
for filename in os.listdir(folder_path):
    if filename.endswith('.pptx'):
        # Create a presentation object for the PowerPoint file
        presentation = Presentation(os.path.join(folder_path, filename))

        # Iterate through the slides in the presentation
        for slide in presentation.slides:
            # Iterate through the shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    # Set the font color to black (RGBColor(0, 0, 0))
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(0, 0, 0)  # Black color
                            run.font.name = "Amasis MT Pro Black"  # Font: Times New Roman
                            run.font.size = Pt(40)  # Font size: 18 points
                            run.font.bold = False  # Remove bold if needed
                            run.font.italic = False  # Remove italic if needed
                            paragraph.alignment = PP_ALIGN.LEFT  # Left-align text

        # Save the modified presentation
        presentation.save(os.path.join(folder_path, filename))

print("Font color and style updated in all PowerPoint files.")
