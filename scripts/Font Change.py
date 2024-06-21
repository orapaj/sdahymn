import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

folder = r"C:\Users\Admin\Desktop\White BG 1"

for filename in os.listdir(folder):
    if filename.endswith('.pptx'):
        prs = Presentation(os.path.join(folder, filename))

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame

                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Amasis MT Pro Black'
                            run.font.size = Pt(39)
                            run.font.color.rgb = RGBColor(0, 0, 0)

        prs.save(os.path.join(folder, filename))

print("Text updated!")
