import os
from pptx import Presentation

folder = r"C:\Users\Admin\Desktop\Output pptx"

for filename in os.listdir(folder):
  if filename.endswith('.pptx'):
    
    prs = Presentation(os.path.join(folder, filename))
    
    for slide in prs.slides:
    
      textboxes = [shp for shp in slide.shapes if shp.has_text_frame]
	  
      max_tb = max(textboxes, key=lambda tb: len(tb.text.split()))

      for tb in textboxes:
        if tb != max_tb:
          tb._element.getparent().remove(tb._element)

    prs.save(filename)

print("Script complete!")
