import os
from pptx import Presentation

def delete_notes_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)

    for slide in presentation.slides:
        # Remove notes from each slide
        slide.notes_slide.notes_text_frame.clear()

    # Save the modified presentation
    presentation.save(pptx_path)

def process_folder(folder_path):
    for filename in os.listdir(folder_path):
        if filename.endswith('.pptx'):
            pptx_path = os.path.join(folder_path, filename)
            delete_notes_from_pptx(pptx_path)

if __name__ == "__main__":
    folder_path = r"D:\Jelmar Orapa\Bible Study Presentations_2\Bible Study Presentations"
    process_folder(folder_path)
