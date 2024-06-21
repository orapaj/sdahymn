import os
import time
import ctypes
import pyautogui
import pythoncom
import subprocess
import win32api
import win32gui
import win32con
import threading
import win32com.client
import tkinter as tk
from tkinter import Scrollbar, Listbox, Entry, Button, Menu, messagebox
from PIL import Image, ImageTk
from tkinter import filedialog
from tkinter import StringVar
from screeninfo import get_monitors
from pptx import Presentation


dir_path = os.path.dirname(os.path.realpath(__file__))
        
def search_files(event=None):
    search_term = search_var.get().lower()
    allowed_extensions = [".pps", ".ppsx"]
    search_results = []

    for root, dirs, files in os.walk(dir_path, topdown=True):
        for file in files:
            if any(file.lower().endswith(ext) for ext in allowed_extensions) and search_term in file.lower():
                search_results.append(file)

    result_listbox.delete(0, tk.END)

    if not search_results:
        result_listbox.insert(tk.END, "No hymn found with that word in the title. Try another!")
    else:
        for result in search_results:
            # Remove the file extension from the result before displaying
            result_without_extension = os.path.splitext(result)[0]
            result_listbox.insert(tk.END, result_without_extension)

def quit_powerpoint_in_thread():
  """Quits PowerPoint in a background thread."""
  try:
    # Initialize COM threading model in the thread
    pythoncom.CoInitialize()

    # Get a reference to the PowerPoint application
    ppt = win32com.client.Dispatch('PowerPoint.Application')

    # Quit the PowerPoint application
    ppt.Quit()
  except Exception as e:
    print("Error:", e)
  finally:
    # Uninitialize COM in the thread
    pythoncom.CoUninitialize()

def quit_powerpoint():
  """Quits PowerPoint in a separate thread to avoid freezing."""
  thread = threading.Thread(target=quit_powerpoint_in_thread)
  thread.start()

def focus_slideshow(event=None):
    time.sleep(0.5)
    screen_width, screen_height = pyautogui.size()
    x = screen_width // 2
    y = 0  
    pyautogui.click(x, y)

def open_selected(event):
    
    selected_item_index = result_listbox.curselection()
    if selected_item_index:
        selected_item = result_listbox.get(selected_item_index)
        selected_file_with_extension = None

    # Recursively search for the selected file in dir_path and its subfolders
    for root, _, files in os.walk(dir_path):
        for file in files:
            if selected_item.lower() in file.lower():
                selected_file_with_extension = os.path.join(root, file)
                break
        if selected_file_with_extension:
            break

    if selected_file_with_extension:
        try:
            # Launch PowerPoint
            ppt = win32com.client.Dispatch('PowerPoint.Application')
            ppt.Visible = True  # Make PowerPoint visible (optional)

            # Open the presentation
            presentation = ppt.Presentations.Open(selected_file_with_extension, WithWindow=True)

            # Switch to Presenter View
            presentation.SlideShowSettings.AdvanceMode = 1  # Advance on click
            presentation.SlideShowSettings.ShowType = 1  # Presenter View
            slide_show = presentation.SlideShowSettings.Run()
            slide_show.Activate()

            # Minimize the PowerPoint window to the taskbar
            ppt.WindowState = 2  # 2 = wdWindowStateMinimized
            focus_slideshow()

        except Exception as e:
            print("Error opening presentation in Presenter View:", e)
            return

def update_background():
    global resized_bg_image
    bg_image = Image.open(r"_internal/Data\bg.png")  # Replace with your image file path
    resized_bg_image = bg_image.resize((root.winfo_width(), root.winfo_height()), Image.LANCZOS)
    bg_image_tk = ImageTk.PhotoImage(resized_bg_image)
    background_label.config(image=bg_image_tk)
    background_label.image = bg_image_tk

def toggle_focus(event=None):
    if search_entry.focus_get() == search_entry:
        result_listbox.select_set(0)
        result_listbox.focus_set()
    else:
        result_listbox.select_clear(0, tk.END)
        search_entry.focus_set()
        search_files()

def clear_search_entry():
    search_entry.delete(0, tk.END)

def select_next_result(event):
    current_selection = result_listbox.curselection()
    if current_selection:
        next_index = (current_selection[0] + 1) % result_listbox.size()
        if next_index == 0:  # Check if the next index is the first item
            next_index = current_selection[0]  # Keep the selection on the current item
        result_listbox.select_clear(current_selection)
        result_listbox.select_set(next_index)
        result_listbox.event_generate("<<ListboxSelect>>")

def select_previous_result(event):
    current_selection = result_listbox.curselection()
    if current_selection:
        previous_index = current_selection[0] - 1
        if previous_index < 0:
            previous_index = 0
        result_listbox.select_clear(current_selection)
        result_listbox.select_set(previous_index)
        result_listbox.event_generate("<<ListboxSelect>>")

def add_hymns():
    file_paths = filedialog.askopenfilenames(
        title="Select Hymn Files",
        filetypes=[("PowerPoint Files", "*.pps *.ppsx")])
    
    if file_paths:
        hymns_directory = os.path.join(dir_path, "Data", "4 More Hymns")
        os.makedirs(hymns_directory, exist_ok=True)

        for file_path in file_paths:
            file_name = os.path.basename(file_path)
            destination_path = os.path.join(hymns_directory, file_name)
            try:
                shutil.copy(file_path, destination_path)
            except Exception as e:
                messagebox.showerror("Error", f"Failed to copy {file_name}: {str(e)}")

        messagebox.showinfo("Success", f"{len(file_paths)} hymn(s) added successfully!")

def helps():
    tk.messagebox.showinfo("Help", "Keyboard Shortcuts: \n\nShift (Right): - Switch between search entry and results' list. \nArrow Up/Down: - Select from the results' list up or down. \nEnter: - To open the selected hymn. \nEsc: - To close or exit from the current hymn played. \n\nAdd Hymns: \n\nTo add hymns that are not on the app's database, \nclick on `File` from the menu bar and select `Add Hymns`, \nthen from the file dialog, select the hymns you want to add. \n\nNote that only .pps or .ppsx file formats are accepted.")
       
def about():
    tk.messagebox.showinfo("About", "Seventh Day Adventist Church Hymnal. \n\nDeveloper: Jelmar A. Orapa \nEmail: orapajelmar@gmail.com")


def open_popup():
    # Calculate the center position of the screen
    screen_width = root.winfo_screenwidth()
    screen_height = root.winfo_screenheight()
    popup_width = 300  # Adjust as needed
    popup_height = 215  # Adjust as needed
    x_position = (screen_width - popup_width) // 2
    y_position = (screen_height - popup_height) // 2

    # Create a new window for the text input
    popup = tk.Toplevel(root)
    popup.title("Announcement!")
    
    popup.geometry(f"{popup_width}x{popup_height}+{x_position}+{y_position}")
    popup.resizable(False, False)

    # Create a Text widget that can accommodate multiple lines
    text_input = tk.Text(popup, wrap=tk.WORD, height=11, width=38)
    text_input.grid(row=0, column=0, sticky=tk.E+tk.W+tk.N+tk.S)
    text_input.focus_set()

    def display_and_close():
        user_input = text_input.get("1.0", "end-1c")  # Get the text entered by the user
        if user_input:
            # Save the text as a .ppsx file
            save_path = save_as_ppsx(user_input)
            if save_path:
                # Open the saved .ppsx file
                os.startfile(save_path)

    # Create an "OK" button to submit the text
    ok_button = tk.Button(popup, text="Display", command=display_and_close, bg="#66A6FF", font=("Times New Roman", 14, "bold"))
    ok_button.grid(row=1, column=0, sticky="nsew")

    def generate_random_filename():
        # Generate a random string of letters and digits
        letters_and_digits = string.ascii_letters + string.digits
        random_string = ''.join(random.choice(letters_and_digits) for i in range(6))

        # Generate a timestamp (current time) to make the filename unique
        timestamp = int(time.time())

        return f"presentation_{timestamp}_{random_string}.pps"

    def save_as_ppsx(text, filename="temp", defaultextension=".pps"):

        prs = Presentation() 
        slide_layout = prs.slide_layouts[0]

        slide = prs.slides.add_slide(slide_layout)
        
        top_title = slide.shapes.title 
        top_title.text = "Announcement!"
        top_title.top = int(prs.slide_height / 10)  # Convert to integer
        top_title.left = 0  # Left align
        top_title.width = prs.slide_width
        top_title.height = int(prs.slide_height / 10)  # Convert to integer


        # Calculate the position and dimensions of the text box
        text_box_left = 0
        text_box_top = top_title.top + top_title.height + (prs.slide_height / 30)  # Below the top title with additional spacing
        text_box_width = prs.slide_width
        text_box_height = prs.slide_height - text_box_top  # Adjust height to fit the remaining slide space


        text_box = slide.shapes.add_textbox(text_box_left, text_box_top, text_box_width, text_box_height)
        text_frame = text_box.text_frame
        p = text_frame.add_paragraph()
        p.text = text


        words = text.split()
        num_words = len(words)
        
        if num_words <= 10:
            font_size = 80 
        elif num_words <= 15:
            font_size = 66
        elif num_words <= 20: 
            font_size = 60
        elif num_words <= 25: 
            font_size = 56
        elif num_words <= 30:
            font_size = 50
        elif num_words <= 40: 
            font_size = 44
        elif num_words <= 50: 
            font_size = 38
        elif num_words <= 75: 
            font_size = 32
        elif num_words <= 100: 
            font_size = 28
        elif num_words <= 150: 
            font_size = 24
        elif num_words <= 200: 
            font_size = 20
        elif num_words <= 250: 
            font_size = 18
        elif num_words <= 300: 
            font_size = 16
        else:
            font_size = 14

        # Set the font size and other properties for the text
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER
        
        # Set text wrapping within the text box
        text_frame.word_wrap = True

        filename = generate_random_filename() if filename == "temp" else filename
        temp_folder = os.path.join(os.getcwd(), "Temp")
        os.makedirs(temp_folder, exist_ok=True)

        save_path = os.path.join(temp_folder, f"{filename}{defaultextension}")

        prs.save(save_path)
        try:
            os.startfile(save_path)
        except Exception as e:
            print(f"Error Displaying Texts: {e}")

def delete_temp_folder():
    try:
        temp_folder = os.path.join(os.getcwd(), "Temp")

        # Check if the "Temp" folder exists
        if os.path.exists(temp_folder):
            # Delete the "Temp" folder and its contents
            shutil.rmtree(temp_folder)
            tk.messagebox.showinfo("Success", "Temp folder deleted successfully!")
        else:
            tk.messagebox.showinfo("Info", "Temp folder does not exist.")

    except Exception as e:
        tk.messagebox.showerror("Error", f"Failed to delete Temp folder: {str(e)}")



# Create the main application window
root = tk.Tk()
root.title("Seventh Day Adventist Church Hymnal")

# Load the image file and convert it to a PhotoImage object
icon_image = Image.open("_internal/Data/favicon.ico")
icon_photo = ImageTk.PhotoImage(icon_image)
root.winfo_toplevel().iconphoto(True, icon_photo)

# Set window dimensions and position
window_width = 510
window_height = 322
screen_width = root.winfo_screenwidth()
screen_height = root.winfo_screenheight()
x_position = (screen_width - window_width) // 2
y_position = (screen_height - window_height) // 2
root.geometry(f"{window_width}x{window_height}+{x_position}+{y_position}")
root.resizable(False, False)

# Create and set the background image
background_label = tk.Label(root)
background_label.place(relwidth=1, relheight=1)
update_background()

# Create a menu bar
menu_bar = Menu(root)
root.config(menu=menu_bar)

# Create a File menu
file_menu = Menu(menu_bar, tearoff=0)
menu_bar.add_cascade(label="File", menu=file_menu)
file_menu.add_command(label="Add Hymns", command=lambda: add_hymns())
file_menu.add_command(label="Delete Temporary Files", command=lambda: delete_temp_folder())
#file_menu.add_separator()
file_menu.add_command(label="Exit", command=root.destroy)

# Create a Help menu
help_menu = Menu(menu_bar, tearoff=0)
menu_bar.add_cascade(label="More", menu=help_menu)
help_menu.add_command(label="Help", command=helps)
help_menu.add_command(label="About", command=about)

# Create a separator in the menu bar
separator = Menu(menu_bar, tearoff=0)
menu_bar.add_cascade(label="I", menu=separator)

#Project a text
menu_bar.add_command(label="Announcement", command=open_popup)

menu_bar.add_command(label="Clean", command=quit_powerpoint)

search_var = StringVar()
# Create a search entry and search button
search_entry = Entry(root, highlightbackground="white", highlightthickness=1, textvariable=search_var)
search_entry.grid(row=0, column=1, padx=0, pady=0)
search_var.trace_add("write", lambda *args: search_files())

search_entry.focus_set()

search_button = Button(root, text="Search", command=search_files)
search_button.grid(row=0, column=2, padx=5, pady=0)

# Create a listbox to display search results
result_listbox = Listbox(root, selectmode=tk.SINGLE, borderwidth=0, highlightthickness=0)
scrollbar = Scrollbar(root, orient=tk.VERTICAL)
scrollbar.config(command=result_listbox.yview)
result_listbox.config(yscrollcommand=scrollbar.set, font=("Times New Roman", 12))
scrollbar.grid(row=1, column=1, padx=0, pady=(0, 24), sticky="ns", rowspan=3)

search_entry.place(in_=result_listbox, x=0, y=0, relx=0.7, relwidth=0.2, relheight=0.07)
search_button.place(in_=result_listbox, x=1, y=0, relx=0.885, relwidth=0.1, relheight=0.07)
search_entry.lift()
search_button.lift()
result_listbox.grid(row=1, column=0, padx=10, pady=(0, 24), sticky="nsew", rowspan=3, columnspan=3)

result_listbox.bind("<Double-Button-1>", open_selected)
result_listbox.bind("<Return>", open_selected)

root.grid_rowconfigure(1, weight=1)
root.grid_columnconfigure(0, weight=1)

root.bind("<Configure>", lambda event: update_background())
root.bind("<Shift_R>", lambda event: toggle_focus())
root.bind("<Up>", select_previous_result)
root.bind("<Down>", select_next_result)

search_files()
root.mainloop()
